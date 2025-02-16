# file_converter.py
import os
from pdf2image import convert_from_path
from ocr import extract_text_from_image
from pptx import Presentation
import shutil
import docx2txt
from PIL import Image

def process_pdf(file_path, poppler_path, progress_callback=None):
    """Convert a PDF to text by converting each page to an image and performing OCR."""
    images = convert_from_path(file_path, dpi=300, poppler_path=poppler_path)
    extracted_text_list = []
    total_pages = len(images)
    for i, image in enumerate(images):
        page_text = extract_text_from_image(image)
        extracted_text_list.append(f"\n\n--- Page {i+1} ---\n\n{page_text}")
        if progress_callback:
            progress_callback(i + 1, total_pages)
    return "".join(extracted_text_list)

def process_image(file_path):
    """Perform OCR directly on an image file."""
    image = Image.open(file_path)
    text = extract_text_from_image(image)
    return text

def process_pptx(file_path):
    """Extract text from a PPTX file using python-pptx."""
    prs = Presentation(file_path)
    extracted_text_list = []
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text += shape.text + "\n"
        extracted_text_list.append(f"\n\n--- Slide {i+1} ---\n\n{slide_text}")
    return "".join(extracted_text_list)

def process_docx(file_path, image_dir='temp_images'):
    """
    Extract text from a DOCX file using docx2txt.
    This function saves all images in the DOCX to a temporary folder,
    then runs OCR on each image and replaces the placeholder in the text
    with the OCR result.
    """
    # Create temporary directory for images if it doesn't exist
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)
    
    # Extract text from the DOCX; docx2txt will also extract images to image_dir.
    text = docx2txt.process(file_path, image_dir)
    
    # docx2txt usually inserts placeholders in the text for images.
    # In many cases the placeholder is of the form "[image:filename.ext]"
    # You can adjust this if your DOCX doesn't produce such placeholders.
    for filename in os.listdir(image_dir):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
            image_path = os.path.join(image_dir, filename)
            # Run OCR on the extracted image
            ocr_text = extract_text_from_image(image_path)
            # Define the placeholder to look for in the extracted text.
            placeholder = f"[image:{filename}]"
            # Replace the placeholder with the OCR result.
            replacement = f"[Image OCR: {ocr_text.strip()}]"
            text = text.replace(placeholder, replacement)
    
    # Optionally, delete the temporary images folder once processing is complete.
    # shutil.rmtree(image_dir)
    
    return text


def file_to_text(file_path, poppler_path, progress_callback=None):
    """
    Convert a file to text. Supports:
      - PDF (.pdf)
      - Images (.jpg, .jpeg, .png, .tiff, .bmp)
      - PowerPoint (.pptx)
      - Word Documents (.docx, .doc)
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return process_pdf(file_path, poppler_path, progress_callback)
    elif ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
        return process_image(file_path)
    elif ext == ".pptx":
        return process_pptx(file_path)
    elif ext in [".docx", ".doc"]:
        return process_docx(file_path)
    else:
        raise ValueError("Unsupported file type")
